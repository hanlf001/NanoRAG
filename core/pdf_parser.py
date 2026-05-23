import os
from pathlib import Path
from typing import List, Dict, Any
import re


class DocumentParser:
    def __init__(self, chunk_size: int = 800, chunk_overlap: int = 100):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.supported_extensions = {
            '.pdf': self._parse_pdf,
            '.txt': self._parse_txt,
            '.md': self._parse_txt,
            '.docx': self._parse_docx,
            '.xlsx': self._parse_xlsx,
            '.pptx': self._parse_pptx
        }

    def parse_document(self, file_path: str) -> Dict[str, Any]:
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        ext = file_path.suffix.lower()
        if ext not in self.supported_extensions:
            raise ValueError(f"不支持的文件格式: {ext}")
        
        parser = self.supported_extensions[ext]
        text = parser(file_path)
        
        chunks = self._split_text(text)
        
        return {
            'file_name': file_path.name,
            'file_path': str(file_path),
            'file_type': ext,
            'total_chars': len(text),
            'chunks': chunks,
            'chunk_count': len(chunks)
        }

    def _split_text(self, text: str) -> List[str]:
        if not text or len(text.strip()) == 0:
            return []
        
        sentences = self._split_into_sentences(text)
        chunks = []
        current_chunk = []
        current_length = 0
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            
            sentence_len = len(sentence)
            
            if current_length + sentence_len > self.chunk_size and current_chunk:
                chunks.append(' '.join(current_chunk))
                overlap = self._calculate_overlap(current_chunk)
                current_chunk = overlap
                current_length = sum(len(s) + 1 for s in current_chunk)
            
            current_chunk.append(sentence)
            current_length += sentence_len + 1
        
        if current_chunk:
            chunks.append(' '.join(current_chunk))
        
        return chunks

    def _split_into_sentences(self, text: str) -> List[str]:
        text = text.replace('\n', ' ')
        text = re.sub(r'\s+', ' ', text)
        
        pattern = r'(?<=[.!?。！？])\s+'
        sentences = re.split(pattern, text)
        sentences = [s.strip() for s in sentences if s.strip()]
        
        return sentences

    def _calculate_overlap(self, chunk: List[str]) -> List[str]:
        if not chunk:
            return []
        
        overlap_length = 0
        overlap = []
        
        for sentence in reversed(chunk):
            if overlap_length + len(sentence) + 1 <= self.chunk_overlap:
                overlap.insert(0, sentence)
                overlap_length += len(sentence) + 1
            else:
                break
        
        return overlap

    def _parse_pdf(self, file_path: Path) -> str:
        import fitz
        doc = fitz.open(file_path)
        text = []
        for page in doc:
            text.append(page.get_text())
        doc.close()
        return '\n'.join(text)

    def _parse_txt(self, file_path: Path) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='gbk') as f:
                return f.read()

    def _parse_docx(self, file_path: Path) -> str:
        from docx import Document
        doc = Document(file_path)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text)
                text.append(' | '.join(row_text))
        return '\n'.join(text)

    def _parse_xlsx(self, file_path: Path) -> str:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True)
        text = []
        for sheet_name in wb.sheetnames:
            text.append(f"Sheet: {sheet_name}")
            sheet = wb[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else '' for cell in row]
                if any(row_text):
                    text.append(' | '.join(row_text))
        wb.close()
        return '\n'.join(text)

    def _parse_pptx(self, file_path: Path) -> str:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            text.append(f"Slide {slide_idx}:")
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text.append(shape.text)
        return '\n'.join(text)


def parse_file(file_path: str, chunk_size: int = 800, chunk_overlap: int = 100) -> Dict[str, Any]:
    parser = DocumentParser(chunk_size, chunk_overlap)
    return parser.parse_document(file_path)
